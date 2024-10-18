import os
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import groq
import tempfile

# Load Groq API key from environment variable
groq_api_key = os.getenv("GROQ_API_KEY")
if not groq_api_key:
    raise ValueError("GROQ_API_KEY environment variable not set")

client = groq.Groq(api_key=groq_api_key)

def generate_slide_content(title, notes):
    prompt = f"Generate detailed content for a presentation slide with the title '{title}'. Use the following notes as a basis:\n{notes}\n\nProvide a structured response with bullet points and sub-points if necessary."
    
    response = client.chat.completions.create(
        messages=[
            {"role": "system", "content": "You are a helpful assistant that generates detailed presentation content."},
            {"role": "user", "content": prompt}
        ],
        model="mixtral-8x7b-32768",
        max_tokens=500
    )
    
    return response.choices[0].message.content

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Using the bullet slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content

def text_to_ppt(text, output_file='output.pptx'):
    prs = Presentation()
    
    # Split the text into paragraphs
    paragraphs = text.split('\n\n')
    
    for para in paragraphs:
        lines = para.split('\n')
        title = lines[0]
        notes = '\n'.join(lines[1:])
        
        # Generate detailed content using Groq
        detailed_content = generate_slide_content(title, notes)
        
        create_slide(prs, title, detailed_content)
    
    prs.save(output_file)
    return output_file

def main():
    st.title("AI-Powered Presentation Generator")
    st.write("Enter your presentation notes below. Each paragraph will be converted into a slide.")

    notes = st.text_area("Enter your notes (separate slides with blank lines):", height=300)
    
    if st.button("Generate Presentation"):
        if notes:
            with st.spinner("Generating presentation..."):
                # Create a temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
                    output_file = tmp_file.name
                    text_to_ppt(notes, output_file)
                
                # Offer the file for download
                with open(output_file, "rb") as file:
                    btn = st.download_button(
                        label="Download Presentation",
                        data=file,
                        file_name="generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                
                # Clean up the temporary file
                os.unlink(output_file)
                
                st.success("Presentation generated successfully!")
        else:
            st.error("Please enter some notes before generating the presentation.")

if __name__ == "__main__":
    main()